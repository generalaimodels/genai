#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
ImageX: A robust, optimized, and generalized Image class with multi-source ingestion, unified constructors,
validation, and an extensible processing pipeline — all in a single Python file.

This file is intentionally verbose with explanatory docstrings and comments so coders can learn deeply.
All examples are runnable (subject to optional dependency availability) and guarded by try/except.

Highlights
- Multi-source ingestion with unified constructors:
  * from_local_path, from_file_uri, from_http_url, from_base64, from_bytes
  * from_numpy, from_pil, from_opencv
  * from_screenshot, from_camera, from_video_frame, from_url_stream
  * from_clipboard
  * from_pdf, from_doc, from_svg
  * from_remote_storage (s3://, gs://, az://)
  * from_zip, from_tar
  * from_tensorflow_tensor, from_torch_tensor
  * from_heic, from_dicom, from_hdf5, from_binary_file
  * from_any (auto-dispatch by type/scheme)
- Validation and normalization with safe dtype/channel handling.
- Extensible processing pipeline with composable processors (Resize, Crop, Rotate, Normalize, etc.).
- Optional verbose logging and rich metadata panels using the 'rich' module.
- Clean API, typed hints, and standards-compliant styling to inspire next generation coders.

Verbose metadata display
- Set environment variable IMG_VERBOSE=1 or pass verbose=True to class methods to get rich meta panels.
- If 'rich' isn't installed, logging gracefully falls back to plain prints or no-ops.

Dependencies
- Core: numpy, pillow (PIL)
- Optional (feature-gated): requests, opencv-python, mss, PyMuPDF (fitz) or pdf2image, python-docx, python-pptx,
  cairosvg, boto3, google-cloud-storage, azure-storage-blob, pillow-heif or pyheif, pydicom, h5py,
  tensorflow, torch, rich.

Note: Methods are implemented to degrade gracefully when optional packages are missing, raising a clear error
with install hints. This design favors robustness and portability.

Examples
- See __main__ at the bottom for runnable snippets demonstrating ingestion and pipeline usage.

Author’s note
- Design favors correctness, clarity, and extensibility. Performance-sensitive sections are implemented with
  vectorized numpy ops or pillow/cv2 where appropriate.
"""

from __future__ import annotations

import base64
import contextlib
import dataclasses
from dataclasses import dataclass, field
import datetime as _dt
import io
import json
import os
import platform
import sys
import tarfile
import zipfile
import hashlib
import warnings
from pathlib import Path
from typing import Any, Callable, Dict, Iterable, List, Mapping, MutableMapping, Optional, Sequence, Tuple, Union
from urllib.parse import urlparse, unquote

# ---------- Optional third-party imports (graceful fallbacks) ----------
try:
    import numpy as np
except Exception as e:
    raise RuntimeError("numpy is required for ImageX.") from e

try:
    from PIL import Image as PILImage, ImageOps as PILImageOps, ImageCms as PILImageCms, ImageGrab as PILImageGrab
except Exception as e:
    raise RuntimeError("Pillow (PIL) is required for ImageX. pip install pillow") from e

# Optional modules (feature-gated)
try:
    import requests
except Exception:
    requests = None

try:
    import cv2
except Exception:
    cv2 = None

try:
    import mss
except Exception:
    mss = None

try:
    import fitz  # PyMuPDF
except Exception:
    fitz = None

try:
    from pdf2image import convert_from_bytes as pdf2img_from_bytes
except Exception:
    pdf2img_from_bytes = None

try:
    import docx  # python-docx
except Exception:
    docx = None

try:
    import pptx  # python-pptx
except Exception:
    pptx = None

try:
    import cairosvg
except Exception:
    cairosvg = None

try:
    import boto3
except Exception:
    boto3 = None

try:
    from google.cloud import storage as gcs_storage
except Exception:
    gcs_storage = None

try:
    from azure.storage.blob import BlobServiceClient
except Exception:
    BlobServiceClient = None

try:
    import h5py
except Exception:
    h5py = None

try:
    import tensorflow as tf
except Exception:
    tf = None

try:
    import torch
except Exception:
    torch = None

try:
    import pillow_heif
    pillow_heif_available = True
    # Register HEIF opener for PIL if available
    try:
        pillow_heif.register_heif_opener()
    except Exception:
        pass
except Exception:
    pillow_heif_available = False
    try:
        import pyheif
    except Exception:
        pyheif = None

try:
    import pydicom
except Exception:
    pydicom = None

# ---------- Rich console (optional) ----------
try:
    from rich.console import Console
    from rich.panel import Panel
    from rich.table import Table
    from rich.text import Text
    from rich.pretty import Pretty
    RICH_AVAILABLE = True
except Exception:
    Console = None
    Panel = None
    Table = None
    Text = None
    Pretty = None
    RICH_AVAILABLE = False

# Verbose toggle via env or per-call override
GLOBAL_VERBOSE: bool = os.getenv("IMG_VERBOSE", "0") in ("1", "true", "True", "YES", "yes")

def _console() -> Optional[Console]:
    if RICH_AVAILABLE:
        return Console()
    return None

def _log_panel(title: str, data: Mapping[str, Any], verbose: Optional[bool]) -> None:
    """Display metadata in a rich panel if verbose is enabled and rich is available."""
    v = GLOBAL_VERBOSE if verbose is None else bool(verbose)
    if not v:
        return
    if RICH_AVAILABLE:
        con = _console()
        if con is None:
            return
        try:
            tbl = Table(show_header=True, header_style="bold magenta")
            tbl.add_column("Key", style="cyan", no_wrap=True)
            tbl.add_column("Value", style="green")
            for k, val in data.items():
                tbl.add_row(str(k), json.dumps(val, default=str) if isinstance(val, (dict, list, tuple)) else str(val))
            con.print(Panel(tbl, title=f"[bold blue]{title}", border_style="blue"))
        except Exception:
            # Fallback to print
            print(f"[{title}]")
            for k, val in data.items():
                print(f"  - {k}: {val}")
    else:
        # Minimal fallback logging
        print(f"[{title}]")
        for k, val in data.items():
            print(f"  - {k}: {val}")

def _hash_bytes(data: bytes, algo: str = "sha256") -> str:
    h = hashlib.new(algo)
    h.update(data)
    return h.hexdigest()

def _now_iso() -> str:
    return _dt.datetime.now().isoformat(timespec="seconds")

def _ensure_uint8(arr: np.ndarray) -> np.ndarray:
    """Convert array to uint8 range safely. Handles 8/16/32-bit and float arrays."""
    if arr.dtype == np.uint8:
        return arr
    a = arr
    if a.dtype.kind in ("f",):
        a = np.nan_to_num(a, nan=0.0, posinf=255.0, neginf=0.0)
        if a.max() <= 1.0:
            a = a * 255.0
        a = np.clip(a, 0, 255).astype(np.uint8)
        return a
    # Integer types
    info = np.iinfo(a.dtype) if a.dtype.kind in ("u", "i") else None
    if info:
        amin, amax = a.min(), a.max()
        if amin == amax:
            if amin <= 0:
                return np.zeros_like(a, dtype=np.uint8)
            return np.full_like(a, 255, dtype=np.uint8)
        # Rescale linearly
        a = (a.astype(np.float32) - amin) / (amax - amin) * 255.0
        a = np.clip(a, 0, 255).astype(np.uint8)
        return a
    # Fallback
    return a.astype(np.uint8)

def _to_hwc(arr: np.ndarray) -> np.ndarray:
    """Ensure HWC layout. Accepts HWC, CHW, HW, HW1, 1HW, etc."""
    if arr.ndim == 2:
        return arr[..., None]  # HW -> HW1
    if arr.ndim == 3:
        h, w, c = arr.shape
        # Heuristics: prefer HWC; if first dim is small (<=4) and last isn't, assume CHW
        if h <= 4 and arr.transpose(1, 2, 0).ndim == 3 and arr.shape[0] <= 4:
            return arr.transpose(1, 2, 0)  # CHW -> HWC
        return arr
    raise ValueError(f"Unsupported array shape for image: {arr.shape}")

def _infer_mode_from_channels(c: int) -> str:
    return {1: "L", 3: "RGB", 4: "RGBA"}.get(c, f"X{c}")

def _pil_to_numpy(pil_img: PILImage.Image, force_uint8: bool = True) -> Tuple[np.ndarray, str]:
    """Convert PIL to numpy HWC, uint8 by default, and return (array, mode)."""
    mode = pil_img.mode
    if mode in ("1", "P"):
        pil_img = pil_img.convert("RGBA")
        mode = "RGBA"
    if mode not in ("L", "RGB", "RGBA", "I;16", "I", "F"):
        try:
            pil_img = pil_img.convert("RGBA")
            mode = "RGBA"
        except Exception:
            pil_img = pil_img.convert("RGB")
            mode = "RGB"
    arr = np.array(pil_img)
    arr = _to_hwc(arr)
    if force_uint8:
        arr = _ensure_uint8(arr)
    c = arr.shape[2] if arr.ndim == 3 else 1
    if c not in (1, 3, 4):
        # Reduce to 3 channels if too many (fallback)
        arr = arr[..., :3]
        mode = "RGB"
    return arr, mode

def _numpy_to_pil(arr: np.ndarray, mode: Optional[str] = None) -> PILImage.Image:
    """Convert numpy HWC to PIL image, inferring mode if not provided."""
    if arr.ndim == 2:
        mode = mode or "L"
        return PILImage.fromarray(_ensure_uint8(arr), mode=mode)
    if arr.ndim != 3:
        raise ValueError("Expected HWC array for PIL conversion.")
    arr = _ensure_uint8(arr)
    h, w, c = arr.shape
    if mode is None:
        mode = _infer_mode_from_channels(c)
    if c == 1 and mode == "L":
        return PILImage.fromarray(arr[..., 0], mode="L")
    if c == 3 and mode in ("RGB", "YCbCr"):
        return PILImage.fromarray(arr, mode="RGB")
    if c == 4 and mode == "RGBA":
        return PILImage.fromarray(arr, mode="RGBA")
    # Fallback: convert to RGB
    if c == 4:
        return PILImage.fromarray(arr[..., :3], mode="RGB")
    if c == 1:
        return PILImage.fromarray(arr[..., 0], mode="L")
    return PILImage.fromarray(arr, mode="RGB")

def _maybe_cv2_to_rgb(arr_bgr: np.ndarray) -> np.ndarray:
    if arr_bgr.ndim == 3 and arr_bgr.shape[2] == 3:
        return arr_bgr[:, :, ::-1]  # BGR->RGB
    if arr_bgr.ndim == 2:
        return arr_bgr[..., None]
    if arr_bgr.ndim == 3 and arr_bgr.shape[2] == 4:
        # BGRA -> RGBA
        return arr_bgr[:, :, [2, 1, 0, 3]]
    return arr_bgr

# ---------- Exceptions ----------
class ImageXError(Exception):
    pass

class DependencyError(ImageXError):
    pass

class ValidationError(ImageXError):
    pass

class SourceError(ImageXError):
    pass

class NotSupportedError(ImageXError):
    pass

# ---------- Processor Pipeline ----------
class Processor:
    """Base processor. Override __call__(self, image: 'Image') -> 'Image'."""
    name: str = "Processor"

    def __call__(self, image: "Image") -> "Image":
        raise NotImplementedError

    def __repr__(self) -> str:
        return f"{self.__class__.__name__}()"

@dataclass
class Pipeline:
    """Composable pipeline of processors."""
    steps: List[Processor] = field(default_factory=list)

    def add(self, *processors: Processor) -> "Pipeline":
        self.steps.extend(processors)
        return self

    def __call__(self, image: "Image") -> "Image":
        for step in self.steps:
            image = step(image)
            image._history.append({
                "op": step.__class__.__name__,
                "params": dataclasses.asdict(step) if dataclasses.is_dataclass(step) else {},
                "timestamp": _now_iso()
            })
        return image

# ---------- Built-in processors ----------
@dataclass
class ToMode(Processor):
    mode: str = "RGB"
    name: str = "ToMode"

    def __call__(self, image: "Image") -> "Image":
        pil = image.to_pil()
        pil = pil.convert(self.mode)
        arr, mode = _pil_to_numpy(pil)
        return image._replace_array(arr, mode=mode, note=f"ToMode({self.mode})")

@dataclass
class Resize(Processor):
    width: Optional[int] = None
    height: Optional[int] = None
    keep_aspect: bool = True
    resample: str = "lanczos"  # nearest, bilinear, bicubic, lanczos
    name: str = "Resize"

    def __call__(self, image: "Image") -> "Image":
        pil = image.to_pil()
        resample_map = {
            "nearest": PILImage.NEAREST,
            "bilinear": PILImage.BILINEAR,
            "bicubic": PILImage.BICUBIC,
            "lanczos": PILImage.LANCZOS,
        }
        rs = resample_map.get(self.resample, PILImage.LANCZOS)
        w, h = pil.size
        if self.width is None and self.height is None:
            return image
        if self.keep_aspect:
            if self.width and not self.height:
                scale = self.width / w
                new_size = (self.width, max(1, int(round(h * scale))))
            elif self.height and not self.width:
                scale = self.height / h
                new_size = (max(1, int(round(w * scale))), self.height)
            else:
                # fit in box
                scale = min(self.width / w, self.height / h)
                new_size = (max(1, int(round(w * scale))), max(1, int(round(h * scale))))
        else:
            new_size = (self.width or w, self.height or h)
        pil = pil.resize(new_size, rs)
        arr, mode = _pil_to_numpy(pil)
        return image._replace_array(arr, mode=mode, note=f"Resize({new_size})")

@dataclass
class CenterCrop(Processor):
    width: int
    height: int
    name: str = "CenterCrop"

    def __call__(self, image: "Image") -> "Image":
        pil = image.to_pil()
        w, h = pil.size
        left = max(0, (w - self.width) // 2)
        top = max(0, (h - self.height) // 2)
        right = min(w, left + self.width)
        bottom = min(h, top + self.height)
        pil = pil.crop((left, top, right, bottom))
        arr, mode = _pil_to_numpy(pil)
        return image._replace_array(arr, mode=mode, note=f"CenterCrop({self.width}x{self.height})")

@dataclass
class Rotate(Processor):
    degrees: float
    expand: bool = True
    fill: Optional[Tuple[int, int, int]] = None
    name: str = "Rotate"

    def __call__(self, image: "Image") -> "Image":
        pil = image.to_pil()
        pil = pil.rotate(self.degrees, expand=self.expand, fillcolor=self.fill)
        arr, mode = _pil_to_numpy(pil)
        return image._replace_array(arr, mode=mode, note=f"Rotate({self.degrees})")

@dataclass
class Flip(Processor):
    horizontal: bool = True
    name: str = "Flip"

    def __call__(self, image: "Image") -> "Image":
        pil = image.to_pil()
        pil = PILImageOps.mirror(pil) if self.horizontal else PILImageOps.flip(pil)
        arr, mode = _pil_to_numpy(pil)
        return image._replace_array(arr, mode=mode, note=f"Flip({'H' if self.horizontal else 'V'})")

@dataclass
class Normalize(Processor):
    mean: Tuple[float, float, float] = (0.485, 0.456, 0.406)
    std: Tuple[float, float, float] = (0.229, 0.224, 0.225)
    inplace: bool = False
    name: str = "Normalize"

    def __call__(self, image: "Image") -> "Image":
        arr = image.array.astype(np.float32) / 255.0
        if arr.ndim == 2:
            arr = arr[..., None]
        c = arr.shape[2]
        m = np.array(self.mean[:c], dtype=np.float32).reshape((1, 1, c))
        s = np.array(self.std[:c], dtype=np.float32).reshape((1, 1, c))
        arr = (arr - m) / s
        # Keep normalized float32 image and tag mode as "F"
        new = image._replace_array(arr, mode="F", note="Normalize")
        return new

@dataclass
class Pad(Processor):
    pad: Tuple[int, int, int, int] = (0, 0, 0, 0)  # left, top, right, bottom
    fill: Tuple[int, int, int] = (0, 0, 0)
    name: str = "Pad"

    def __call__(self, image: "Image") -> "Image":
        pil = image.to_pil()
        l, t, r, b = self.pad
        nw, nh = pil.size[0] + l + r, pil.size[1] + t + b
        background = PILImage.new("RGB", (nw, nh), self.fill)
        background.paste(pil, (l, t))
        arr, mode = _pil_to_numpy(background)
        return image._replace_array(arr, mode=mode, note=f"Pad({self.pad},{self.fill})")

# ---------- Image class ----------
@dataclass
class Image:
    """
    An immutable image object with multi-source ingestion, validation,
    and an extensible processing pipeline.

    Core representation
    - array: numpy ndarray (H, W, C) or (H, W) for grayscale. Prefer uint8 for visual tasks.
    - mode: string like 'RGB', 'RGBA', 'L', or 'F' for float arrays.
    - metadata: dict capturing source, uri/path, creation time, and arbitrary extra info.
    - history: list of operations applied via pipeline processors.

    Immutability
    - All operations return new Image instances. Internal helpers _replace_array produce clones.

    Validation
    - validate() enforces non-empty image, numeric dtype, shape compatibility, and sane channels.

    Conversion
    - to_pil, to_numpy, to_opencv, to_base64, to_bytes, save, show.

    Unified constructors cover many sources. Each classmethod optionally logs metadata via rich (IMG_VERBOSE=1).

    Examples:
    - Loading from various sources:
        img = Image.from_local_path("photo.jpg", verbose=True)
        img = Image.from_http_url("https://picsum.photos/256", verbose=True)
        img = Image.from_numpy(np.zeros((128,128,3), dtype=np.uint8))
        img = Image.from_pil(PILImage.open("photo.png"))

    - Processing pipeline:
        pipeline = Pipeline().add(ToMode("RGB"), Resize(width=512), CenterCrop(512, 512), Normalize())
        out = pipeline(img)
        out.save("out.png")

    - Auto-dispatch:
        out = Image.from_any("file:///path/to/img.png")
        out = Image.from_any(b"...raw bytes...")
        out = Image.from_any(np.zeros((32,32,3), dtype=np.uint8))
    """
    array: np.ndarray
    mode: str = "RGB"
    metadata: Dict[str, Any] = field(default_factory=dict)
    _history: List[Dict[str, Any]] = field(default_factory=list, repr=False)

    # ------------------------ Core internal helpers ------------------------
    def _clone(self) -> "Image":
        return Image(array=self.array.copy(), mode=self.mode, metadata=self.metadata.copy(), _history=self._history.copy())

    def _replace_array(self, new_array: np.ndarray, mode: Optional[str] = None, note: Optional[str] = None) -> "Image":
        img = Image(array=new_array, mode=mode or self.mode, metadata=self.metadata.copy(), _history=self._history.copy())
        if note:
            img._history.append({"op": "replace_array", "note": note, "timestamp": _now_iso()})
        return img

    # ---------------------------- Validation ------------------------------
    def validate(self) -> "Image":
        if not isinstance(self.array, np.ndarray):
            raise ValidationError("array must be numpy.ndarray")
        if self.array.size == 0:
            raise ValidationError("Image array is empty.")
        if self.array.ndim not in (2, 3):
            raise ValidationError(f"Unsupported array ndim={self.array.ndim}, expected 2 or 3.")
        if self.array.ndim == 3 and self.array.shape[2] not in (1, 2, 3, 4):
            raise ValidationError("Channels must be 1, 2, 3, or 4.")
        if self.mode not in ("L", "RGB", "RGBA", "F", "I", "YCbCr", "CMYK", "LAB", "HSV"):
            # Not strictly enforced, but warn
            warnings.warn(f"Unknown mode '{self.mode}'. Proceeding anyway.")
        return self

    # ----------------------------- Converters -----------------------------
    def to_pil(self) -> PILImage.Image:
        """Convert to Pillow Image; float images converted reasonably."""
        arr = self.array
        if arr.dtype != np.uint8 and self.mode != "F":
            arr = _ensure_uint8(arr)
        if self.mode == "F" and arr.dtype != np.float32:
            arr = arr.astype(np.float32)
        return _numpy_to_pil(arr, mode=None if self.mode == "F" else self.mode)

    def to_numpy(self, copy: bool = True) -> np.ndarray:
        """Return underlying array, optionally copying for safety."""
        return self.array.copy() if copy else self.array

    def to_opencv(self) -> np.ndarray:
        """Return BGR/BGRA array compatible with OpenCV interface."""
        arr = self.array
        if arr.ndim == 2:
            return arr
        if arr.shape[2] == 3:
            return arr[:, :, ::-1]  # RGB->BGR
        if arr.shape[2] == 4:
            return arr[:, :, [2, 1, 0, 3]]  # RGBA->BGRA
        return arr

    def to_bytes(self, format: str = "PNG", **pil_save_kwargs: Any) -> bytes:
        """Encode to bytes via PIL."""
        pil = self.to_pil()
        with io.BytesIO() as buf:
            pil.save(buf, format=format, **pil_save_kwargs)
            return buf.getvalue()

    def to_base64(self, format: str = "PNG", urlsafe: bool = False, **pil_save_kwargs: Any) -> str:
        """Encode to base64 string (ascii)."""
        raw = self.to_bytes(format=format, **pil_save_kwargs)
        b = base64.urlsafe_b64encode(raw) if urlsafe else base64.b64encode(raw)
        return b.decode("ascii")

    # ------------------------------- IO -----------------------------------
    def save(self, path: Union[str, os.PathLike], format: Optional[str] = None, makedirs: bool = True, **pil_save_kwargs: Any) -> None:
        """Save image to path. Automatically creates directories if needed."""
        path = Path(path)
        if makedirs:
            path.parent.mkdir(parents=True, exist_ok=True)
        pil = self.to_pil()
        pil.save(str(path), format=format, **pil_save_kwargs)

    def show(self, title: Optional[str] = None) -> None:
        """Display image using PIL's default viewer."""
        pil = self.to_pil()
        pil.show(title=title)

    # --------------------------- Info / Metadata ---------------------------
    @property
    def width(self) -> int:
        return int(self.array.shape[1])

    @property
    def height(self) -> int:
        return int(self.array.shape[0])

    @property
    def channels(self) -> int:
        return 1 if self.array.ndim == 2 else int(self.array.shape[2])

    def info(self, verbose: Optional[bool] = None) -> Dict[str, Any]:
        meta = {
            "mode": self.mode,
            "shape": list(self.array.shape),
            "dtype": str(self.array.dtype),
            "width": self.width,
            "height": self.height,
            "channels": self.channels,
            "metadata": self.metadata,
            "history": self._history[-5:],  # last 5 ops
        }
        _log_panel("Image.info", meta, verbose)
        return meta

    # ----------------------------- Pipeline --------------------------------
    def process(self, pipeline: Union[Pipeline, Iterable[Processor]], verbose: Optional[bool] = None) -> "Image":
        """Apply pipeline; returns a new Image."""
        if not isinstance(pipeline, Pipeline):
            pipeline = Pipeline(list(pipeline))
        result = pipeline(self)
        _log_panel("Image.process", {"steps": [s.__class__.__name__ for s in pipeline.steps], "out_shape": list(result.array.shape)}, verbose)
        return result

    # ------------------------ Unified constructors -------------------------

    @classmethod
    def _finalize(cls, arr: np.ndarray, mode: Optional[str], metadata: Dict[str, Any], verbose: Optional[bool]) -> "Image":
        arr = _to_hwc(arr) if arr.ndim != 2 else arr
        if mode is None:
            c = 1 if arr.ndim == 2 else arr.shape[2]
            mode = _infer_mode_from_channels(c)
        img = cls(array=arr, mode=mode, metadata=metadata, _history=[{"op": "create", "timestamp": _now_iso(), "meta": metadata.copy()}])
        img.validate()
        _log_panel("Image created", {"shape": list(img.array.shape), "mode": img.mode, "meta": metadata}, verbose)
        return img

    # --- Path/URI/URL ---
    @classmethod
    def from_local_path(cls, path: str, *, verbose: Optional[bool] = None) -> "Image":
        """Load image from local filesystem path (absolute or relative)."""
        p = Path(path).expanduser()
        if not p.exists():
            raise SourceError(f"File not found: {p}")
        pil = PILImage.open(str(p))
        arr, mode = _pil_to_numpy(pil)
        meta = {"source": "local_path", "path": str(p.resolve()), "format": pil.format, "created": _now_iso()}
        return cls._finalize(arr, mode, meta, verbose)

    @classmethod
    def from_file_uri(cls, uri: str, *, verbose: Optional[bool] = None) -> "Image":
        """Load from file:// URI."""
        parsed = urlparse(uri)
        if parsed.scheme != "file":
            raise SourceError(f"Not a file URI: {uri}")
        path = unquote(parsed.path)
        return cls.from_local_path(path, verbose=verbose)

    @classmethod
    def from_http_url(cls, url: str, *, timeout: int = 20, headers: Optional[Dict[str, str]] = None, verbose: Optional[bool] = None) -> "Image":
        """Download and decode image from http/https."""
        if requests is None:
            raise DependencyError("requests is required. pip install requests")
        r = requests.get(url, headers=headers, timeout=timeout)
        r.raise_for_status()
        data = r.content
        pil = PILImage.open(io.BytesIO(data))
        arr, mode = _pil_to_numpy(pil)
        meta = {"source": "http_url", "url": url, "status": r.status_code, "content_type": r.headers.get("Content-Type"), "hash": _hash_bytes(data), "created": _now_iso()}
        return cls._finalize(arr, mode, meta, verbose)

    @classmethod
    def from_url_stream(cls, url: str, *, api: str = "auto", verbose: Optional[bool] = None) -> "Image":
        """
        Read a single frame from an IP camera or streaming URL.
        api: 'opencv'|'auto'. 'opencv' uses cv2.VideoCapture which supports many streams (e.g., MJPEG).
        """
        if api in ("opencv", "auto"):
            if cv2 is None:
                if api == "opencv":
                    raise DependencyError("OpenCV required for from_url_stream. pip install opencv-python")
            if cv2 is not None:
                cap = cv2.VideoCapture(url)
                if not cap.isOpened():
                    raise SourceError(f"Cannot open stream: {url}")
                ok, frame = cap.read()
                cap.release()
                if not ok:
                    raise SourceError("Failed to read frame from stream.")
                frame = _maybe_cv2_to_rgb(frame)
                frame = _ensure_uint8(frame)
                meta = {"source": "url_stream", "url": url, "api": "opencv", "created": _now_iso()}
                return cls._finalize(frame, None, meta, verbose)
        raise NotSupportedError("No suitable backend found to read from URL stream.")

    # --- Memory/Encodings ---
    @classmethod
    def from_base64(cls, data: str, *, verbose: Optional[bool] = None) -> "Image":
        """Decode base64-encoded image string into Image."""
        raw = base64.b64decode(data)
        return cls.from_bytes(raw, verbose=verbose, base64_input=True)

    @classmethod
    def from_bytes(cls, data: bytes, *, verbose: Optional[bool] = None, base64_input: bool = False) -> "Image":
        """Construct image from raw byte buffer."""
        try:
            pil = PILImage.open(io.BytesIO(data))
        except Exception as e:
            raise SourceError(f"Unable to decode bytes as image: {e}")
        arr, mode = _pil_to_numpy(pil)
        meta = {"source": "bytes", "size_bytes": len(data), "format": pil.format, "hash": _hash_bytes(data), "from_base64": base64_input, "created": _now_iso()}
        return cls._finalize(arr, mode, meta, verbose)

    # --- In-memory tensors/arrays ---
    @classmethod
    def from_numpy(cls, array: np.ndarray, *, mode: Optional[str] = None, copy: bool = False, verbose: Optional[bool] = None) -> "Image":
        """Accept in-memory NumPy array as image tensor; infers mode from channels if not provided."""
        arr = np.array(array, copy=copy)
        arr = _to_hwc(arr) if arr.ndim != 2 else arr
        if arr.ndim == 3 and arr.shape[2] not in (1, 3, 4):
            # Reduce/extract first 3 channels for safety
            arr = arr[..., :3]
            mode = mode or "RGB"
        if arr.dtype != np.uint8 and (mode != "F"):
            arr = _ensure_uint8(arr)
        c = 1 if arr.ndim == 2 else arr.shape[2]
        mode = mode or _infer_mode_from_channels(c)
        meta = {"source": "numpy", "shape": list(arr.shape), "dtype": str(arr.dtype), "created": _now_iso()}
        return cls._finalize(arr, mode, meta, verbose)

    @classmethod
    def from_pil(cls, pil_img: PILImage.Image, *, verbose: Optional[bool] = None) -> "Image":
        arr, mode = _pil_to_numpy(pil_img)
        meta = {"source": "pil", "size": pil_img.size, "format": pil_img.format, "created": _now_iso()}
        return cls._finalize(arr, mode, meta, verbose)

    @classmethod
    def from_opencv(cls, cv2_img: np.ndarray, *, verbose: Optional[bool] = None) -> "Image":
        """Use OpenCV-captured frame or matrix as source."""
        arr = _maybe_cv2_to_rgb(cv2_img)
        arr = _ensure_uint8(arr)
        meta = {"source": "opencv", "shape": list(arr.shape), "created": _now_iso()}
        return cls._finalize(arr, None, meta, verbose)

    # --- OS / Device ---
    @classmethod
    def from_screenshot(cls, os_name: Optional[str] = None, *, monitor: int = 1, region: Optional[Tuple[int, int, int, int]] = None, verbose: Optional[bool] = None) -> "Image":
        """
        Capture screen using mss (cross-platform).
        region: (left, top, width, height) if provided.
        """
        if mss is None:
            raise DependencyError("mss required for screenshots. pip install mss")
        os_name = os_name or platform.system()
        with mss.mss() as sct:
            mon = sct.monitors[monitor] if monitor < len(sct.monitors) else sct.monitors[0]
            if region is not None:
                left, top, w, h = region
                bbox = {"left": left, "top": top, "width": w, "height": h}
            else:
                bbox = mon
            shot = sct.grab(bbox)
            arr = np.array(shot)
            # mss returns BGRA
            arr = arr[..., :3][:, :, ::-1]  # -> RGB
            meta = {"source": "screenshot", "os": os_name, "region": region, "monitor": monitor, "created": _now_iso()}
            return cls._finalize(arr, "RGB", meta, verbose)

    @classmethod
    def from_camera(cls, device_index: int = 0, *, verbose: Optional[bool] = None) -> "Image":
        """Acquire a single frame from system webcam/USB camera."""
        if cv2 is None:
            raise DependencyError("OpenCV required. pip install opencv-python")
        cap = cv2.VideoCapture(device_index)
        if not cap.isOpened():
            raise SourceError(f"Cannot open camera index {device_index}")
        ok, frame = cap.read()
        cap.release()
        if not ok:
            raise SourceError("Failed to capture from camera.")
        frame = _maybe_cv2_to_rgb(frame)
        frame = _ensure_uint8(frame)
        meta = {"source": "camera", "device_index": device_index, "created": _now_iso()}
        return cls._finalize(frame, None, meta, verbose)

    @classmethod
    def from_video_frame(cls, video_path: str, frame_no: int = 0, *, verbose: Optional[bool] = None) -> "Image":
        """Extract still image from video using OpenCV."""
        if cv2 is None:
            raise DependencyError("OpenCV required. pip install opencv-python")
        cap = cv2.VideoCapture(str(video_path))
        if not cap.isOpened():
            raise SourceError(f"Cannot open video: {video_path}")
        cap.set(cv2.CAP_PROP_POS_FRAMES, frame_no)
        ok, frame = cap.read()
        cap.release()
        if not ok:
            raise SourceError(f"Failed to read frame {frame_no} from video.")
        frame = _maybe_cv2_to_rgb(frame)
        frame = _ensure_uint8(frame)
        meta = {"source": "video_frame", "path": str(video_path), "frame_no": frame_no, "created": _now_iso()}
        return cls._finalize(frame, None, meta, verbose)

    @classmethod
    def from_clipboard(cls, *, verbose: Optional[bool] = None) -> "Image":
        """Read image from system clipboard (supported on Windows/macOS; limited on some Linux)."""
        try:
            pil = PILImageGrab.grabclipboard()
        except Exception as e:
            raise SourceError(f"Clipboard access failed: {e}")
        if pil is None or not isinstance(pil, PILImage.Image):
            raise SourceError("No image found in clipboard.")
        arr, mode = _pil_to_numpy(pil)
        meta = {"source": "clipboard", "created": _now_iso()}
        return cls._finalize(arr, mode, meta, verbose)

    # --- Documents & Vector ---
    @classmethod
    def from_pdf(cls, pdf_path: str, page_no: int = 0, *, dpi: int = 200, verbose: Optional[bool] = None) -> "Image":
        """Render a PDF page to raster image using PyMuPDF or pdf2image."""
        p = Path(pdf_path).expanduser()
        if not p.exists():
            raise SourceError(f"PDF not found: {p}")
        data = p.read_bytes()
        if fitz is not None:
            doc = fitz.open(stream=data, filetype="pdf")
            if page_no < 0 or page_no >= doc.page_count:
                raise SourceError(f"Page {page_no} out of range 0..{doc.page_count-1}")
            page = doc.load_page(page_no)
            mat = fitz.Matrix(dpi / 72, dpi / 72)
            pix = page.get_pixmap(matrix=mat)
            arr = np.frombuffer(pix.samples, dtype=np.uint8)
            if pix.alpha:
                arr = arr.reshape((pix.height, pix.width, 4))
                mode = "RGBA"
            else:
                arr = arr.reshape((pix.height, pix.width, 3))
                mode = "RGB"
            meta = {"source": "pdf", "path": str(p.resolve()), "page": page_no, "dpi": dpi, "created": _now_iso()}
            return cls._finalize(arr, mode, meta, verbose)
        elif pdf2img_from_bytes is not None:
            pil_pages = pdf2img_from_bytes(data, dpi=dpi, first_page=page_no + 1, last_page=page_no + 1)
            pil = pil_pages[0]
            arr, mode = _pil_to_numpy(pil)
            meta = {"source": "pdf", "path": str(p.resolve()), "page": page_no, "dpi": dpi, "backend": "pdf2image", "created": _now_iso()}
            return cls._finalize(arr, mode, meta, verbose)
        else:
            raise DependencyError("Install PyMuPDF (pip install pymupdf) or pdf2image (plus poppler) to render PDFs.")

    @classmethod
    def from_doc(cls, document_path: str, *, index: int = 0, verbose: Optional[bool] = None) -> "Image":
        """
        Extract an embedded image from document formats.
        Supported: .docx (python-docx), .pptx (python-pptx).
        """
        p = Path(document_path).expanduser()
        if not p.exists():
            raise SourceError(f"Document not found: {p}")
        suffix = p.suffix.lower()
        if suffix == ".docx":
            if docx is None:
                raise DependencyError("python-docx required. pip install python-docx")
            d = docx.Document(str(p))
            # Embedded images live in part.package.parts with content_type like image/*
            imgs = [part for part in d.part.package.parts if part.content_type and part.content_type.startswith("image/")]
            if not imgs or index >= len(imgs):
                raise SourceError("No embedded images found in DOCX or index out of range.")
            data = imgs[index].blob
            return cls.from_bytes(data, verbose=verbose)
        elif suffix == ".pptx":
            if pptx is None:
                raise DependencyError("python-pptx required. pip install python-pptx")
            pres = pptx.Presentation(str(p))
            # Collect pictures from slides
            images: List[bytes] = []
            for slide in pres.slides:
                for shape in slide.shapes:
                    with contextlib.suppress(Exception):
                        if getattr(shape, "shape_type", None) and shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                            images.append(shape.image.blob)
            if not images or index >= len(images):
                raise SourceError("No embedded images found in PPTX or index out of range.")
            return cls.from_bytes(images[index], verbose=verbose)
        else:
            raise NotSupportedError("Only .docx and .pptx are supported in this single-file implementation.")

    @classmethod
    def from_svg(cls, svg_path: str, *, output_dpi: int = 96, verbose: Optional[bool] = None) -> "Image":
        """Rasterize SVG to PNG using cairosvg."""
        if cairosvg is None:
            raise DependencyError("cairosvg required to rasterize SVG. pip install cairosvg")
        p = Path(svg_path).expanduser()
        if not p.exists():
            raise SourceError(f"SVG not found: {p}")
        png_bytes = cairosvg.svg2png(url=str(p.resolve()), dpi=output_dpi)
        return cls.from_bytes(png_bytes, verbose=verbose)

    # --- Cloud/Archives ---
    @classmethod
    def from_remote_storage(cls, uri: str, *, verbose: Optional[bool] = None) -> "Image":
        """
        Fetch image from cloud backends.
        Supported schemes: s3://bucket/key, gs://bucket/blob, az://container/blob
        """
        parsed = urlparse(uri)
        scheme = parsed.scheme.lower()
        if scheme == "s3":
            if boto3 is None:
                raise DependencyError("boto3 required for S3. pip install boto3")
            bucket = parsed.netloc
            key = parsed.path.lstrip("/")
            s3 = boto3.client("s3")
            obj = s3.get_object(Bucket=bucket, Key=key)
            data = obj["Body"].read()
            return cls.from_bytes(data, verbose=verbose)
        elif scheme == "gs":
            if gcs_storage is None:
                raise DependencyError("google-cloud-storage required. pip install google-cloud-storage")
            client = gcs_storage.Client()
            bucket = client.bucket(parsed.netloc)
            blob = bucket.blob(parsed.path.lstrip("/"))
            data = blob.download_as_bytes()
            return cls.from_bytes(data, verbose=verbose)
        elif scheme in ("az", "azure", "azureblob"):
            if BlobServiceClient is None:
                raise DependencyError("azure-storage-blob required. pip install azure-storage-blob")
            # Expect az://account/container/blob or az://container/blob with env AZ_CONNECTION_STRING
            conn_str = os.getenv("AZ_CONNECTION_STRING")
            if not conn_str:
                raise DependencyError("Set AZ_CONNECTION_STRING for Azure Blob.")
            bsc = BlobServiceClient.from_connection_string(conn_str)
            parts = parsed.path.strip("/").split("/", 1)
            if parsed.netloc and parts:
                container = parsed.netloc
                blob_name = parts[0] if len(parts) == 1 else "/".join(parts)
            elif len(parts) == 2:
                container, blob_name = parts
            else:
                raise SourceError("Invalid Azure Blob URI. Use az://container/blob_path")
            blob_client = bsc.get_blob_client(container=container, blob=blob_name)
            data = blob_client.download_blob().readall()
            return cls.from_bytes(data, verbose=verbose)
        else:
            raise NotSupportedError(f"Unsupported remote scheme: {scheme}")

    @classmethod
    def from_zip(cls, archive_path: str, filename: str, *, verbose: Optional[bool] = None) -> "Image":
        """Read image directly inside .zip archive without extraction."""
        p = Path(archive_path).expanduser()
        if not p.exists():
            raise SourceError(f"Zip file not found: {p}")
        with zipfile.ZipFile(str(p), "r") as zf:
            with zf.open(filename, "r") as fp:
                data = fp.read()
        return cls.from_bytes(data, verbose=verbose)

    @classmethod
    def from_tar(cls, tar_path: str, member: str, *, verbose: Optional[bool] = None) -> "Image":
        """Extract and decode image from tar/tar.gz without full extraction."""
        p = Path(tar_path).expanduser()
        if not p.exists():
            raise SourceError(f"Tar file not found: {p}")
        mode = "r:gz" if p.suffixes[-1:] == [".gz"] or p.suffix == ".tgz" else "r:*"
        with tarfile.open(str(p), mode) as tfp:
            m = tfp.getmember(member)
            f = tfp.extractfile(m)
            if f is None:
                raise SourceError(f"Member not found: {member}")
            data = f.read()
        return cls.from_bytes(data, verbose=verbose)

    # --- Framework tensors ---
    @classmethod
    def from_tensorflow_tensor(cls, tf_tensor: Any, *, verbose: Optional[bool] = None) -> "Image":
        """
        Convert TensorFlow tensor to image.
        Accepts BHWC/HWC/CHW. Takes the first item if batch is present.
        """
        if tf is None:
            raise DependencyError("tensorflow required. pip install tensorflow")
        arr = tf_tensor.numpy() if hasattr(tf_tensor, "numpy") else np.array(tf_tensor)
        if arr.ndim == 4:
            arr = arr[0]
        arr = _to_hwc(arr)
        if arr.dtype != np.uint8:
            arr = _ensure_uint8(arr)
        meta = {"source": "tensorflow", "shape": list(arr.shape), "dtype": str(arr.dtype), "created": _now_iso()}
        return cls._finalize(arr, None, meta, verbose)

    @classmethod
    def from_torch_tensor(cls, torch_tensor: Any, *, verbose: Optional[bool] = None) -> "Image":
        """
        Convert PyTorch tensor to image.
        Accepts BCHW/CHW/HWC. Takes the first item if batch is present.
        """
        if torch is None:
            raise DependencyError("torch required. pip install torch")
        t = torch_tensor
        with torch.no_grad():
            arr = t.detach().cpu().numpy()
        if arr.ndim == 4:
            arr = arr[0]
        arr = _to_hwc(arr)
        if arr.dtype != np.uint8:
            arr = _ensure_uint8(arr)
        meta = {"source": "torch", "shape": list(arr.shape), "dtype": str(arr.dtype), "created": _now_iso()}
        return cls._finalize(arr, None, meta, verbose)

    # --- Formats ---
    @classmethod
    def from_heic(cls, heic_path: str, *, verbose: Optional[bool] = None) -> "Image":
        """Decode HEIC/HEIF images via pillow-heif (preferred) or pyheif."""
        p = Path(heic_path).expanduser()
        if not p.exists():
            raise SourceError(f"HEIC not found: {p}")
        if pillow_heif_available:
            pil = PILImage.open(str(p))
            arr, mode = _pil_to_numpy(pil)
            meta = {"source": "heic", "path": str(p.resolve()), "backend": "pillow-heif", "created": _now_iso()}
            return cls._finalize(arr, mode, meta, verbose)
        elif 'pyheif' in globals() and pyheif is not None:
            heif = pyheif.read(str(p))
            img = PILImage.frombytes(heif.mode, heif.size, heif.data, "raw", heif.mode, heif.stride)
            arr, mode = _pil_to_numpy(img)
            meta = {"source": "heic", "path": str(p.resolve()), "backend": "pyheif", "created": _now_iso()}
            return cls._finalize(arr, mode, meta, verbose)
        else:
            raise DependencyError("Install pillow-heif or pyheif to decode HEIC.")

    @classmethod
    def from_dicom(cls, dicom_path: str, *, window: Optional[Tuple[int, int]] = None, verbose: Optional[bool] = None) -> "Image":
        """Parse and render DICOM using pydicom with windowing to 8-bit RGB/Grayscale."""
        if pydicom is None:
            raise DependencyError("pydicom required. pip install pydicom")
        p = Path(dicom_path).expanduser()
        if not p.exists():
            raise SourceError(f"DICOM not found: {p}")
        ds = pydicom.dcmread(str(p))
        arr = ds.pixel_array.astype(np.float32)
        # Window/level handling
        if window:
            center = (window[0] + window[1]) / 2.0
            width = (window[1] - window[0])
        else:
            center = float(ds.get("WindowCenter", np.mean(arr)))
            width = float(ds.get("WindowWidth", max(1.0, arr.max() - arr.min())))
        low = center - width / 2.0
        high = center + width / 2.0
        arr = (arr - low) / (high - low)
        arr = np.clip(arr, 0, 1) * 255.0
        arr = arr.astype(np.uint8)
        if arr.ndim == 3:
            if arr.shape[2] == 3:
                mode = "RGB"
            else:
                arr = arr[..., 0]
                mode = "L"
        else:
            mode = "L"
        meta = {"source": "dicom", "path": str(p.resolve()), "window": window or [low, high], "created": _now_iso()}
        return cls._finalize(arr, mode, meta, verbose)

    @classmethod
    def from_hdf5(cls, hdf_path: str, dataset: str, *, verbose: Optional[bool] = None) -> "Image":
        """Load image dataset from HDF5 file using h5py."""
        if h5py is None:
            raise DependencyError("h5py required. pip install h5py")
        p = Path(hdf_path).expanduser()
        if not p.exists():
            raise SourceError(f"HDF5 not found: {p}")
        with h5py.File(str(p), "r") as f:
            if dataset not in f:
                raise SourceError(f"Dataset not found in HDF5: {dataset}")
            arr = f[dataset][()]
        arr = _to_hwc(arr) if arr.ndim != 2 else arr
        arr = _ensure_uint8(arr)
        meta = {"source": "hdf5", "path": str(p.resolve()), "dataset": dataset, "shape": list(arr.shape), "created": _now_iso()}
        return cls._finalize(arr, None, meta, verbose)

    @classmethod
    def from_binary_file(cls, bin_path: str, shape: Tuple[int, int, int] | Tuple[int, int], *, dtype: Union[str, np.dtype] = "uint8", verbose: Optional[bool] = None) -> "Image":
        """Interpret raw binary file as image buffer with given shape and dtype."""
        p = Path(bin_path).expanduser()
        if not p.exists():
            raise SourceError(f"Binary file not found: {p}")
        arr = np.fromfile(str(p), dtype=np.dtype(dtype))
        expected = int(np.prod(shape))
        if arr.size < expected:
            raise SourceError(f"Binary file too small. Expected {expected} elements, got {arr.size}.")
        arr = arr[:expected].reshape(shape)
        arr = _to_hwc(arr) if arr.ndim != 2 else arr
        if arr.dtype != np.uint8:
            arr = _ensure_uint8(arr)
        meta = {"source": "binary_file", "path": str(p.resolve()), "shape": list(arr.shape), "dtype": "uint8", "created": _now_iso()}
        return cls._finalize(arr, None, meta, verbose)

    # --- General dispatcher ---
    @classmethod
    def from_any(cls, x: Any, *, verbose: Optional[bool] = None, **kwargs: Any) -> "Image":
        """
        Auto-dispatch based on input type or string scheme.
        Supported:
        - str: http(s)://, file://, s3://, gs://, az://, or local path
        - bytes: raw image bytes
        - numpy.ndarray: array
        - PIL.Image.Image: Pillow image
        - cv2 image (ndarray)
        - torch/tensorflow tensors
        """
        if isinstance(x, str):
            parsed = urlparse(x)
            if parsed.scheme in ("http", "https"):
                return cls.from_http_url(x, verbose=verbose)
            if parsed.scheme == "file":
                return cls.from_file_uri(x, verbose=verbose)
            if parsed.scheme in ("s3", "gs", "az", "azure", "azureblob"):
                return cls.from_remote_storage(x, verbose=verbose)
            # fallback local path
            p = Path(x)
            if p.exists():
                # HEIC special-case
                if p.suffix.lower() in (".heic", ".heif"):
                    return cls.from_heic(x, verbose=verbose)
                if p.suffix.lower() == ".pdf":
                    return cls.from_pdf(x, verbose=verbose)
                if p.suffix.lower() in (".docx", ".pptx"):
                    return cls.from_doc(x, verbose=verbose)
                if p.suffix.lower() == ".svg":
                    return cls.from_svg(x, verbose=verbose)
                return cls.from_local_path(x, verbose=verbose)
            # Base64 heuristic
            with contextlib.suppress(Exception):
                return cls.from_base64(x, verbose=verbose)
            raise SourceError(f"Cannot resolve input string: {x}")
        if isinstance(x, (bytes, bytearray, memoryview)):
            return cls.from_bytes(bytes(x), verbose=verbose)
        if isinstance(x, np.ndarray):
            return cls.from_numpy(x, verbose=verbose, **kwargs)
        if isinstance(x, PILImage.Image):
            return cls.from_pil(x, verbose=verbose)
        # TF tensor
        if tf is not None and hasattr(tf, "is_tensor") and tf.is_tensor(x):
            return cls.from_tensorflow_tensor(x, verbose=verbose)
        # Torch tensor
        if torch is not None and isinstance(x, torch.Tensor):
            return cls.from_torch_tensor(x, verbose=verbose)
        raise NotSupportedError(f"Unsupported input type for from_any: {type(x)}")

# ------------------------ Extended examples and usage ------------------------
def _example_usage() -> None:
    """
    Run a series of self-contained examples. Safe for environments without external resources.
    Toggle metadata verbosity via IMG_VERBOSE=1 or set verbose=True in calls.
    """
    v = True  # make examples verbose for demonstration
    examples: List[Tuple[str, Callable[[], Image]]] = []

    # 1. From NumPy (grayscale and RGB)
    def ex_numpy_rgb() -> Image:
        arr = np.zeros((128, 256, 3), dtype=np.uint8)
        arr[:, :128, 0] = 255  # left half red
        return Image.from_numpy(arr, verbose=v)

    def ex_numpy_gray() -> Image:
        arr = np.linspace(0, 255, 128 * 128, dtype=np.uint8).reshape(128, 128)
        return Image.from_numpy(arr, verbose=v)

    examples.append(("from_numpy RGB", ex_numpy_rgb))
    examples.append(("from_numpy Grayscale", ex_numpy_gray))

    # 2. From bytes/base64
    def ex_bytes_b64() -> Image:
        # Create a tiny image in memory
        tiny = Image.from_numpy(np.full((10, 10, 3), 200, dtype=np.uint8), verbose=v)
        b64 = tiny.to_base64(format="PNG")
        return Image.from_base64(b64, verbose=v)

    examples.append(("from_base64", ex_bytes_b64))

    # 3. PIL wrapper
    def ex_pil() -> Image:
        pil = PILImage.new("RGB", (64, 64), (0, 128, 255))
        return Image.from_pil(pil, verbose=v)

    examples.append(("from_pil", ex_pil))

    # 4. Processing pipeline
    def ex_pipeline() -> Image:
        img = ex_numpy_rgb()
        pipe = Pipeline().add(ToMode("RGB"), Resize(width=64), CenterCrop(64, 64), Rotate(15), Flip(True))
        out = img.process(pipe, verbose=v)
        return out

    examples.append(("pipeline basic", ex_pipeline))

    # 5. Save/Load roundtrip in memory
    def ex_save_load() -> Image:
        img = ex_numpy_gray()
        buf = img.to_bytes(format="PNG")
        return Image.from_bytes(buf, verbose=v)

    examples.append(("save/load bytes", ex_save_load))

    # 6. DICOM (skips if pydicom not available)
    def ex_dicom() -> Image:
        if pydicom is None:
            raise DependencyError("Skipping DICOM example: pydicom missing.")
        raise DependencyError("Provide a real DICOM path to test this example.")
    # Not appending dicom example by default to avoid errors

    # 7. PDF (skips if no backend)
    def ex_pdf() -> Image:
        if fitz is None and pdf2img_from_bytes is None:
            raise DependencyError("Skipping PDF example: no backend.")
        # Generate a simple in-memory PDF via reportlab or fallback: create an SVG and rasterize
        if cairosvg is None:
            raise DependencyError("Install cairosvg to demonstrate PDF/SVG flow in-memory.")
        svg = '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100"><rect x="10" y="10" width="180" height="80" style="fill: #0af; stroke: #000; stroke-width: 3;"/><text x="20" y="60" font-size="24" fill="white">Hello PDF</text></svg>'
        png = cairosvg.svg2png(bytestring=svg.encode("utf-8"))
        return Image.from_bytes(png, verbose=v)
    # Not appending; environment-dependent

    # Run
    for title, fn in examples:
        try:
            img = fn()
            meta = img.info(verbose=v)
            # Attempt a small processing
            out = img.process([Resize(width=32, keep_aspect=True), ToMode("RGB")], verbose=v)
            # Save to temporary path
            out_path = Path("output") / f"{title.replace(' ', '_')}.png"
            out.save(out_path)
            if RICH_AVAILABLE:
                _console().print(Panel.fit(Text(f"Example '{title}' succeeded -> {out_path}", style="bold green"), title="OK"))
            else:
                print(f"[OK] {title} -> {out_path}")
        except Exception as e:
            if RICH_AVAILABLE:
                _console().print(Panel.fit(Text(f"Example '{title}' failed: {e}", style="bold red"), title="ERROR"))
            else:
                print(f"[ERROR] {title}: {e}")

# ------------------------ Notes on exceptions and edge cases ------------------------
"""
Edge cases and safeguards:
- Non-8-bit arrays (float/16-bit) are safely converted to uint8 via _ensure_uint8 for display pipelines.
- Unknown channel counts are trimmed to 3 where possible to avoid downstream errors.
- DICOM windowing is approximated if values aren't provided; medical workflows may need precise VOI LUT handling.
- Clipboard/screenshot/video/camera support may vary by OS, permissions, or drivers.
- from_url_stream uses OpenCV VideoCapture which supports many MJPEG/RTSP streams; for other formats implement custom readers.
- from_remote_storage assumes proper credentials/environment for each cloud provider.
- from_doc supports DOCX/PPTX; legacy binary formats are not handled in this single-file edition.

Extensibility:
- Add new sources as @classmethod on Image and call _finalize(arr, mode, metadata, verbose).
- Add new processors by subclassing Processor and implementing __call__.
- Compose pipelines using Pipeline().add(...). Each step is recorded in history.
"""

# ------------------------ Main: quick demo ------------------------
if __name__ == "__main__":
    # Ensure output directory exists
    Path("output").mkdir(parents=True, exist_ok=True)

    # You can also set IMG_VERBOSE=1 in env to get rich meta panels
    _example_usage()

    # Bonus: auto-dispatch demo
    try:
        # Create an in-memory tiny image as bytes
        tmp = Image.from_numpy(np.full((8, 8, 3), (10, 200, 30), dtype=np.uint8), verbose=True)
        any1 = Image.from_any(tmp.to_bytes(format="PNG"), verbose=True)
        any1.save("output/from_any_bytes.png")
        if RICH_AVAILABLE:
            _console().print(Panel.fit(Text("from_any(bytes) -> output/from_any_bytes.png", style="bold green"), title="from_any"))
        else:
            print("from_any(bytes) -> output/from_any_bytes.png")
    except Exception as e:
        if RICH_AVAILABLE:
            _console().print(Panel.fit(Text(f"from_any demo failed: {e}", style="bold red"), title="ERROR"))
        else:
            print(f"[ERROR] from_any demo: {e}")